"""Business Review service — gather analytics data + render to HTML/PDF/PPTX.

Pipeline:
    gather_data(account_id, cadence, dates) -> snapshot dict
    render_html(snapshot)                   -> str (single-file HTML doc)
    render_pdf(html)                        -> bytes (weasyprint)
    render_pptx(snapshot)                   -> bytes (python-pptx, 12 slides)

The snapshot mirrors the slide structure of
`/Users/anandkaliappan/Desktop/Beroe/Beroe_BR_Mondel_z_International_monthly_2026-06-05.pptx`
1:1. Slides 9/10/11 carry placeholder text in Phase 1; the chart-image
bytes get injected in Phase 3 (matplotlib).

Cadence drives the period window only — the deck structure is identical
across monthly / quarterly / renewal / custom.
"""

from __future__ import annotations

import base64
import html as _html
import logging
from dataclasses import dataclass
from datetime import date, timedelta
from decimal import Decimal
from io import BytesIO
from typing import Any

from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from app.models.account import Account
from app.models.cs_goal import CSGoal
from app.models.metric import SuccessMetric
from app.models.play import AccountPlay
from app.models.signal import SoftSignal

logger = logging.getLogger(__name__)


# ============================================================
# Cadence → period window
# ============================================================


@dataclass(frozen=True)
class PeriodWindow:
    label: str
    start: date | None
    end: date | None


_MONTHS = ("Jan", "Feb", "Mar", "Apr", "May", "Jun",
           "Jul", "Aug", "Sep", "Oct", "Nov", "Dec")


def derive_period(
    *,
    cadence: str,
    today: date,
    contract_start: date | None,
    renewal_date: date | None,
    custom_start: date | None,
    custom_end: date | None,
    label_override: str | None,
) -> PeriodWindow:
    if cadence == "custom":
        if not (custom_start and custom_end):
            raise ValueError("custom cadence requires period_start + period_end")
        label = label_override or f"{custom_start.strftime('%-d %b %Y')} – {custom_end.strftime('%-d %b %Y')}"
        return PeriodWindow(label=label, start=custom_start, end=custom_end)

    if cadence == "monthly":
        start = (today.replace(day=1) - timedelta(days=1)).replace(day=1)
        end = today
        label = label_override or f"{_MONTHS[start.month - 1]} {start.year}"
        return PeriodWindow(label=label, start=start, end=end)

    if cadence == "quarterly":
        q = (today.month - 1) // 3 + 1
        # snap to current quarter start
        q_start = date(today.year, (q - 1) * 3 + 1, 1)
        label = label_override or f"Q{q} {today.year}"
        return PeriodWindow(label=label, start=q_start, end=today)

    if cadence == "renewal":
        start = contract_start or today.replace(month=1, day=1)
        end = renewal_date or today
        label = label_override or (
            f"Renewal review · {end.strftime('%-d %b %Y')}" if renewal_date else "Renewal review"
        )
        return PeriodWindow(label=label, start=start, end=end)

    raise ValueError(f"unknown cadence: {cadence}")


# ============================================================
# Data gather — pulls from every relevant table
# ============================================================


async def gather_data(
    *,
    db: AsyncSession,
    account: Account,
    period: PeriodWindow,
    cadence: str,
) -> dict[str, Any]:
    """Build the 12-section snapshot. Missing fields become em-dashes
    at render time; this function focuses on shape + plumbing."""

    pi = account.platform_intel or {}
    usage = pi.get("usage", {}) or {}
    modules = pi.get("modules", {}) or {}
    cat_intel = pi.get("cat_intel", {}) or {}
    abi = pi.get("abi", {}) or {}
    scores = pi.get("scores", {}) or {}
    cust = pi.get("custom_credits", {}) or {}
    inflation = pi.get("inflation_watch", {}) or {}
    nps = pi.get("nps", {}) or {}

    # — soft signals (risks / open / asks come from M27)
    # signal_status enum: 'active' | 'resolved'
    # signal_impact enum: 'critical' | 'high' | 'medium' | 'low'
    signals_q = await db.execute(
        select(SoftSignal).where(
            SoftSignal.account_id == account.id,
            SoftSignal.status == "active",
            SoftSignal.hidden.is_(False),
        )
    )
    signals = signals_q.scalars().all()
    open_items = [s.signal for s in signals if s.impact in ("critical", "high")]
    risks = [s.signal for s in signals if s.impact == "critical"]
    asks: list[str] = []  # populated below from checkpoints / VDD if present

    # — plays (upsell pipeline)
    plays_q = await db.execute(
        select(AccountPlay).where(
            AccountPlay.account_id == account.id,
            AccountPlay.hidden.is_(False),
        )
    )
    plays = plays_q.scalars().all()
    upsell_pipeline = [
        {
            "title": p.title,
            "status": (p.role or "").upper() or "ACTIVE",
            "value_usd": float(p.value_usd or 0),
            "prob": int(p.prob or 0),
            "when_text": p.when_text or "",
        }
        for p in plays
        if (p.prob or 0) >= 40
    ]

    # — metrics (cost avoidance / value)
    metrics_q = await db.execute(
        select(SuccessMetric).where(
            SuccessMetric.account_id == account.id,
            SuccessMetric.deleted_at.is_(None),
        )
    )
    metrics = metrics_q.scalars().all()
    cost_avoid = 0.0
    cost_target = 0.0
    for m in metrics:
        try:
            cv = _num(m.current_value)
            tv = _num(m.target_value)
            if "savings" in (m.name or "").lower() or "cost" in (m.name or "").lower():
                cost_avoid += cv
                cost_target += tv
        except (TypeError, ValueError):
            pass

    # — goals (accomplishments / asks)
    goals_q = await db.execute(
        select(CSGoal).where(
            CSGoal.account_id == account.id,
            CSGoal.deleted_at.is_(None),
        )
    )
    goals = goals_q.scalars().all()
    accomplishments: list[dict[str, Any]] = []
    for g in goals:
        # Each goal's initiatives in "delivered" stage are wins.
        for init in g.initiatives or []:
            if isinstance(init, dict) and (init.get("stage") in ("delivered", "in_progress")):
                accomplishments.append({
                    "amount_usd": init.get("value_usd") or init.get("amount") or 0,
                    "title": init.get("title") or g.title,
                    "bu": init.get("bu") or g.category,
                    "description": init.get("description") or "",
                })
    # Asks: open initiatives in identification/pipeline
    for g in goals:
        for init in g.initiatives or []:
            if isinstance(init, dict) and init.get("stage") in ("identification", "pipeline"):
                ask = init.get("title")
                if ask:
                    asks.append(ask)

    # — sponsor (from cs_stakeholders)
    stakeholders = account.cs_stakeholders or {}
    sponsor_obj = stakeholders.get("champion") or stakeholders.get("commercial") or {}

    # — derived: days to renewal
    days_to_renewal: int | None = None
    if account.gate_renewal_date:
        days_to_renewal = (account.gate_renewal_date - date.today()).days

    # — scores 12 (from platform_intel.scores, else compute approximations)
    scores_12 = _build_scores_12(scores)

    return {
        "meta": {
            "account_name": account.name,
            "cadence": cadence,
            "period_label": period.label,
            "period_start": period.start.isoformat() if period.start else None,
            "period_end": period.end.isoformat() if period.end else None,
        },
        "cover": {
            "account_name": account.name,
            "period_label": period.label,
            "title": "COMPREHENSIVE BUSINESS REVIEW",
            "footer": "Beroe",
        },
        "custom_credits": {
            "fte": cust.get("fte"),
            "hours_purchased": cust.get("hours_purchased"),
            "hours_consumed": cust.get("hours_consumed"),
            "hours_remaining": cust.get("hours_remaining"),
            "fixed_fee_done": cust.get("fixed_fee_done"),
            "fixed_fee_total": cust.get("fixed_fee_total"),
            "instances_in_flight": cust.get("instances_in_flight"),
            "slots_active": cust.get("slots_active"),
            "slot_cadence": cust.get("slot_cadence"),
            "slots": cust.get("slots", []),
        },
        "risks_open_asks": {
            "open_items": open_items,
            "risks": risks or [s.title for s in signals if s.impact == "risk"],
            "asks": asks[:10],
        },
        "executive_snapshot": {
            "health": scores.get("composite") or scores.get("account_health"),
            "health_trend": scores.get("trend_text"),
            "health_band": scores.get("band"),
            "arr_usd": float(account.current_acv or 0),
            "expansion_note": scores.get("expansion_note") or "",
            "cost_avoidance_usd": cost_avoid,
            "cost_target_usd": cost_target,
            "renewal_date": account.gate_renewal_date.isoformat() if account.gate_renewal_date else None,
            "days_to_renewal": days_to_renewal,
            "active_30d": usage.get("active_30d") or usage.get("active_users"),
            "licensed_seats": usage.get("licensed_seats"),
            "logins": usage.get("logins_period") or usage.get("logins_q1"),
            "hours": usage.get("hours_period") or usage.get("hours_q1"),
            "abi_queries": abi.get("queries_period") or abi.get("queries_q1"),
            "value_text": scores.get("value_text"),
        },
        "contract_summary": {
            "start": account.contract_start.isoformat() if account.contract_start else None,
            "end": account.gate_renewal_date.isoformat() if account.gate_renewal_date else (
                account.contract_end.isoformat() if account.contract_end else None
            ),
            "days_to_renewal": days_to_renewal,
            "term": account.gate_contract_term,
            "acv_usd": float(account.gate_contract_acv or account.current_acv or 0),
            "billing": (pi.get("billing") or {}).get("frequency") or "Annual",
            "seats_current": usage.get("licensed_seats"),
            "seats_proposed": usage.get("seats_proposed"),
            "geography": account.country,
            "sponsor_name": sponsor_obj.get("name"),
            "sponsor_role": sponsor_obj.get("role"),
            "modules": list(account.gate_contract_modules or []),
        },
        "accomplishments": accomplishments[:6],
        "upsell_pipeline": sorted(
            upsell_pipeline, key=lambda x: x["value_usd"], reverse=True
        )[:8],
        "scores_12": scores_12,
        "subscribers_engagement": {
            "licensed_seats": usage.get("licensed_seats"),
            "seats_proposed": usage.get("seats_proposed"),
            "active_30d": usage.get("active_30d"),
            "activation_pct": _pct(
                usage.get("active_30d"), usage.get("licensed_seats")
            ),
            "logins_trend": usage.get("monthly_logins") or [],
            "logins_total": usage.get("logins_period"),
            "hours_total": usage.get("hours_period"),
        },
        "live_ai": {
            "subscribers": cat_intel.get("subscribers"),
            "total_subs": cat_intel.get("total_subs") or usage.get("licensed_seats"),
            "categories_unlocked": cat_intel.get("categories_unlocked"),
            "ent_cats": cat_intel.get("ent_cats"),
            "non_ent_cats": cat_intel.get("non_ent_cats"),
            "avg_per_user": cat_intel.get("avg_per_user"),
            "benchmark": cat_intel.get("benchmark") or 1.8,
            "top_cats": [
                {"name": c.get("name"), "visits": c.get("visits", 0), "heat": c.get("heat", "cold")}
                for c in (cat_intel.get("top_cats") or [])[:10]
            ],
        },
        "inflation_watch": {
            "categories_tracked": inflation.get("categories_tracked"),
            "categories_in_scope": inflation.get("categories_in_scope"),
            "views_period": inflation.get("views_period") or inflation.get("views_q1"),
            "neg_prep_runs": inflation.get("neg_prep_runs"),
            "trend_monthly": inflation.get("trend_monthly") or [],
        },
        "nps": {
            "score": nps.get("score"),
            "promoters": nps.get("promoters"),
            "passives": nps.get("passives"),
            "detractors": nps.get("detractors"),
        },
    }


def _num(v: Any) -> float:
    if v is None:
        return 0.0
    if isinstance(v, (int, float, Decimal)):
        return float(v)
    s = str(v).strip().replace("$", "").replace(",", "").replace("%", "")
    m_unit = s.lower().endswith("m")
    k_unit = s.lower().endswith("k")
    s = s.rstrip("MmKk").strip()
    try:
        n = float(s)
        if m_unit:
            n *= 1_000_000
        elif k_unit:
            n *= 1_000
        return n
    except ValueError:
        return 0.0


def _pct(num: Any, denom: Any) -> int | None:
    try:
        n = float(num or 0)
        d = float(denom or 0)
        if d == 0:
            return None
        return round(100 * n / d)
    except (TypeError, ValueError):
        return None


_SCORE_LABELS = (
    "Account Health (composite)",
    "Product Score",
    "Engagement Score",
    "Adoption Score",
    "Value Score",
    "Renewal Risk Score",
    "Expansion Score",
    "Support Score",
    "NPS Score",
    "Stakeholder Score",
    "Time-to-Value Score",
    "Data Quality Score",
)


def _build_scores_12(scores_block: dict) -> list[dict[str, Any]]:
    """Map platform_intel.scores → 12-row deck table. Always returns 12
    rows; uses em-dash for missing values at render time."""

    rows: list[dict[str, Any]] = []
    score_keys = (
        "composite",
        "product",
        "engagement",
        "adoption",
        "value",
        "renewal_risk",
        "expansion",
        "support",
        "nps",
        "stakeholder",
        "time_to_value",
        "data_quality",
    )
    for label, key in zip(_SCORE_LABELS, score_keys, strict=True):
        s = scores_block.get(key)
        band = "HEALTHY" if s and s >= 70 else "AT RISK" if s and s >= 40 else "CRITICAL" if s else None
        rows.append({"name": label, "score": s, "band": band})
    return rows


# ============================================================
# HTML renderer
# ============================================================


def _e(v: Any) -> str:
    """HTML-escape + em-dash for missing."""
    if v is None or v == "":
        return "—"
    return _html.escape(str(v))


def _money(v: Any) -> str:
    n = _num(v)
    if n == 0:
        return "—"
    if abs(n) >= 1_000_000:
        return f"${n/1_000_000:.1f}M"
    if abs(n) >= 1_000:
        return f"${n/1_000:.0f}K"
    return f"${n:,.0f}"


def _kpi_tile(label: str, value: str, note: str = "") -> str:
    return (
        f'<div class="tile">'
        f'<div class="tile-label">{_e(label)}</div>'
        f'<div class="tile-value">{value}</div>'
        f'<div class="tile-note">{_e(note) if note else ""}</div>'
        f"</div>"
    )


# ============================================================
# Chart rendering — matplotlib PNG bytes, embedded inline in
# HTML (data: URL) AND in PPTX (add_picture). Same bytes →
# identical visuals across all three outputs.
# ============================================================


_INDIGO_HEX = "#4A00F8"
_NAVY_HEX = "#001137"
_FUSCIA_HEX = "#C344C7"
_AQUA_HEX = "#35E1D4"


def _chart_to_png(fig: Any) -> bytes:
    """Convert a matplotlib Figure to PNG bytes, then close it.

    Always closes the figure so the worker process doesn't leak figures
    across calls (matplotlib keeps a global registry by default)."""
    buf = BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", dpi=120,
                facecolor="white", edgecolor="none")
    import matplotlib.pyplot as plt  # type: ignore
    plt.close(fig)
    return buf.getvalue()


def _line_chart(
    *, months: list[str], series: list[float], title: str, color: str = _INDIGO_HEX
) -> bytes:
    """12-month line chart with a soft fill under it. Returns empty
    bytes when the data has fewer than 2 points (a line of 1 point is
    just a dot — render placeholder text instead)."""
    if not months or not series or len(series) < 2:
        return b""
    try:
        import matplotlib  # type: ignore
        matplotlib.use("Agg")  # headless safe
        import matplotlib.pyplot as plt  # type: ignore
    except ImportError:
        logger.warning("matplotlib not installed; skipping line chart")
        return b""

    fig, ax = plt.subplots(figsize=(9.5, 3.2))
    n = min(len(months), len(series))
    x = list(range(n))
    y = list(series[:n])
    labels = list(months[:n])
    ax.plot(x, y, marker="o", linewidth=2.5, color=color, markersize=6,
            markerfacecolor="white", markeredgewidth=2, markeredgecolor=color)
    ax.fill_between(x, y, alpha=0.10, color=color)
    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=9, color="#5a7896")
    ax.tick_params(axis="y", labelsize=9, colors="#5a7896")
    ax.set_title(title, fontsize=11, fontweight="bold", color=_NAVY_HEX, loc="left", pad=12)
    for spine in ("top", "right"):
        ax.spines[spine].set_visible(False)
    ax.spines["left"].set_color("#e4eaf6")
    ax.spines["bottom"].set_color("#e4eaf6")
    ax.grid(axis="y", color="#f0f3f8", linewidth=1)
    ax.set_axisbelow(True)
    return _chart_to_png(fig)


def _bar_chart(
    *, labels: list[str], values: list[float], title: str, color: str = _INDIGO_HEX
) -> bytes:
    """Horizontal bar chart. Used for slide 10's top-categories list."""
    if not labels or not values:
        return b""
    try:
        import matplotlib  # type: ignore
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt  # type: ignore
    except ImportError:
        logger.warning("matplotlib not installed; skipping bar chart")
        return b""

    n = min(len(labels), len(values), 10)
    ll = list(reversed(labels[:n]))
    vv = list(reversed(values[:n]))
    fig, ax = plt.subplots(figsize=(9.5, max(2.4, 0.34 * n + 1.2)))
    bars = ax.barh(ll, vv, color=color, edgecolor="white", linewidth=1)
    for bar, v in zip(bars, vv, strict=True):
        ax.text(v + max(vv) * 0.01, bar.get_y() + bar.get_height() / 2,
                f"{int(v) if float(v).is_integer() else v}",
                va="center", fontsize=9, color=_NAVY_HEX, fontweight="bold")
    ax.tick_params(axis="x", labelsize=9, colors="#5a7896")
    ax.tick_params(axis="y", labelsize=9, colors=_NAVY_HEX)
    ax.set_title(title, fontsize=11, fontweight="bold", color=_NAVY_HEX, loc="left", pad=10)
    for spine in ("top", "right", "bottom"):
        ax.spines[spine].set_visible(False)
    ax.spines["left"].set_color("#e4eaf6")
    ax.set_xticks([])
    return _chart_to_png(fig)


def _build_chart_pngs(snapshot: dict[str, Any]) -> dict[str, bytes]:
    """Build the three chart PNGs once; reuse them across HTML+PPTX.

    Keys:
      slide9  — 12-month login trend (line)
      slide10 — Top categories visits (horizontal bar)
      slide11 — Monthly inflation-watch trend (line)
    """
    se = snapshot.get("subscribers_engagement", {}) or {}
    la = snapshot.get("live_ai", {}) or {}
    iw = snapshot.get("inflation_watch", {}) or {}

    # Slide 9 — login trend. Months come from platform_intel.usage.months
    # but subscribers_engagement.logins_trend mirrors that list. If the
    # mirror is absent fall back to a 12-month default label set.
    months = snapshot.get("meta", {}).get("logins_months") or _MONTHS_12()
    logins = se.get("logins_trend") or []
    chart9 = _line_chart(
        months=months, series=[float(x) for x in logins if isinstance(x, (int, float))],
        title="Monthly logins (12-month trend)",
    )

    # Slide 10 — top categories
    cats = la.get("top_cats") or []
    chart10 = _bar_chart(
        labels=[str(c.get("name") or "—") for c in cats[:10]],
        values=[float(c.get("visits") or 0) for c in cats[:10]],
        title="Top categories by visits",
        color=_FUSCIA_HEX,
    )

    # Slide 11 — inflation trend
    inf_trend = iw.get("trend_monthly") or []
    chart11 = _line_chart(
        months=_MONTHS_12()[: len(inf_trend)],
        series=[float(v) for v in inf_trend if isinstance(v, (int, float))],
        title="Inflation Watch · monthly trend",
        color=_AQUA_HEX,
    )

    return {"slide9": chart9, "slide10": chart10, "slide11": chart11}


def _MONTHS_12() -> list[str]:
    return list(_MONTHS)


def _png_data_url(png: bytes) -> str:
    """Inline a PNG into an HTML <img src="data:..."> URL."""
    if not png:
        return ""
    b64 = base64.b64encode(png).decode("ascii")
    return f"data:image/png;base64,{b64}"


_CSS = """
*{box-sizing:border-box;margin:0;padding:0}
body{font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',Manrope,sans-serif;
     background:#fff;color:#0d1b2e;font-size:13px}
.slide{padding:32px 48px;min-height:560px;page-break-after:always;
       background:#fff;border-bottom:2px solid #e4eaf6}
.slide.cover{background:linear-gradient(135deg,#001137 0%,#0f1d54 100%);
             color:#fff;min-height:720px;display:flex;flex-direction:column;
             justify-content:center;align-items:center;text-align:center}
.slide.cover h1{font-size:42px;font-weight:800;margin-bottom:16px}
.slide.cover .meta{font-size:18px;color:#a3b6e0;letter-spacing:2px;margin-bottom:8px}
.slide.cover .title{font-size:24px;font-weight:700;margin:24px 0 8px;color:#FFE61E}
.slide.cover .brand{font-size:14px;color:#a3b6e0;margin-top:32px}
.slide-eyebrow{font-size:11px;font-weight:700;letter-spacing:2px;color:#5a7896;
               margin-bottom:8px;text-transform:uppercase}
.slide-title{font-size:24px;font-weight:800;color:#001137;margin-bottom:24px;
             padding-bottom:12px;border-bottom:2px solid #4A00F8}
.grid{display:grid;gap:16px}
.grid-3{grid-template-columns:repeat(3,1fr)}
.grid-4{grid-template-columns:repeat(4,1fr)}
.grid-2{grid-template-columns:repeat(2,1fr)}
.tile{padding:16px;border:1px solid #e4eaf6;border-radius:10px;
      background:linear-gradient(135deg,#fff 0%,#f7f9fd 100%)}
.tile-label{font-size:10px;font-weight:700;color:#5a7896;letter-spacing:1.5px;
            text-transform:uppercase;margin-bottom:6px}
.tile-value{font-size:24px;font-weight:800;color:#001137;line-height:1.1}
.tile-note{font-size:11px;color:#5a7896;margin-top:6px}
table{width:100%;border-collapse:collapse;font-size:12px}
th{text-align:left;padding:8px 10px;background:#f7f9fd;color:#5a7896;
   font-size:10px;font-weight:700;letter-spacing:1px;text-transform:uppercase}
td{padding:10px;border-bottom:1px solid #f0f3f8}
.pill{display:inline-block;padding:3px 8px;border-radius:6px;font-size:10px;
      font-weight:700;letter-spacing:0.5px}
.pill.green{background:#dcfce7;color:#146a45}
.pill.amber{background:#fef3c7;color:#854F0B}
.pill.red{background:#fee2e2;color:#8B1F1F}
.pill.blue{background:#dbeafe;color:#1e40af}
.bullet{padding:8px 0;border-bottom:1px dashed #e4eaf6;display:flex;gap:10px}
.bullet:before{content:"•";color:#4A00F8;font-weight:bold}
h3.section{font-size:13px;font-weight:700;color:#001137;
           margin:20px 0 10px;letter-spacing:0.5px}
.dollar-card{padding:18px;border-left:3px solid #4A00F8;background:#f7f9fd;
             border-radius:8px;margin-bottom:10px}
.dollar-card .amt{font-size:22px;font-weight:800;color:#001137}
.dollar-card .label{font-size:12px;font-weight:700;color:#5a7896;margin-top:4px}
.dollar-card .desc{font-size:11px;color:#5a7896;margin-top:6px}
"""


def render_html(snapshot: dict[str, Any], *, charts: dict[str, bytes] | None = None) -> str:
    s = snapshot
    cover = s["cover"]
    cc = s["custom_credits"]
    ro = s["risks_open_asks"]
    ex = s["executive_snapshot"]
    cs = s["contract_summary"]
    acc = s.get("accomplishments") or []
    up = s.get("upsell_pipeline") or []
    sc = s.get("scores_12") or []
    se = s["subscribers_engagement"]
    la = s["live_ai"]
    iw = s["inflation_watch"]
    # Charts: build once if caller didn't pass them. Empty bytes means
    # the data was too thin to chart and we fall back to a text note.
    if charts is None:
        charts = _build_chart_pngs(snapshot)

    def slide(eyebrow: str, title: str, body: str) -> str:
        eb = f'<div class="slide-eyebrow">·  {_e(eyebrow)}</div>' if eyebrow else ""
        return f'<section class="slide">{eb}<h2 class="slide-title">{_e(title)}</h2>{body}</section>'

    # Slide 1 — Cover
    s1 = f"""<section class="slide cover">
      <div class="meta">{_e(cover.get('period_label'))}</div>
      <h1>{_e(cover.get('account_name'))}</h1>
      <div class="title">{_e(cover.get('title'))}</div>
      <div class="brand">{_e(cover.get('footer'))}</div>
    </section>"""

    # Slide 2 — Custom Credits
    slot_rows = "".join(
        f'<tr><td>{_e(x.get("module"))}</td><td>{_e(x.get("owner"))}</td></tr>'
        for x in (cc.get("slots") or [])
    ) or '<tr><td colspan="2">—</td></tr>'
    s2_body = f"""
      <div class="grid grid-3">
        {_kpi_tile("FTE allocation", _e(cc.get('fte')))}
        {_kpi_tile("Hours purchased", _e(cc.get('hours_purchased')))}
        {_kpi_tile("Hours consumed YTD", _e(cc.get('hours_consumed')))}
        {_kpi_tile("Hours remaining", _e(cc.get('hours_remaining')))}
        {_kpi_tile("Fixed-fee instances", f"{_e(cc.get('fixed_fee_done'))} / {_e(cc.get('fixed_fee_total'))}")}
        {_kpi_tile("Instances in flight", _e(cc.get('instances_in_flight')))}
        {_kpi_tile("Infinity Slots active", _e(cc.get('slots_active')))}
        {_kpi_tile("Slot refresh cadence", _e(cc.get('slot_cadence')))}
      </div>
      <h3 class="section">Active slots</h3>
      <table><thead><tr><th>Module</th><th>Owner</th></tr></thead>
        <tbody>{slot_rows}</tbody></table>"""

    # Slide 3 — Risks / Open / Asks
    def _bullets(items: list[Any]) -> str:
        if not items:
            return '<div class="bullet">—</div>'
        return "".join(f'<div class="bullet">{_e(x)}</div>' for x in items)
    s3_body = f"""
      <div class="grid grid-3">
        <div><h3 class="section">Open items (Beroe-side)</h3>{_bullets(ro.get('open_items') or [])}</div>
        <div><h3 class="section">Risks</h3>{_bullets(ro.get('risks') or [])}</div>
        <div><h3 class="section">Asks from client</h3>{_bullets(ro.get('asks') or [])}</div>
      </div>"""

    # Slide 4 — Executive Snapshot
    s4_body = f"""
      <div class="grid grid-4">
        {_kpi_tile("ACCOUNT HEALTH", f"{_e(ex.get('health'))} / 100", _e(ex.get('health_trend') or ex.get('health_band')))}
        {_kpi_tile("ARR IN SCOPE", _money(ex.get('arr_usd')), _e(ex.get('expansion_note')))}
        {_kpi_tile("COST AVOIDANCE", _money(ex.get('cost_avoidance_usd')), f"target {_money(ex.get('cost_target_usd'))}" if ex.get('cost_target_usd') else "")}
        {_kpi_tile("RENEWAL", _e(ex.get('renewal_date')), f"{_e(ex.get('days_to_renewal'))} days" if ex.get('days_to_renewal') is not None else "")}
      </div>
      <h3 class="section">Engagement</h3>
      <div class="grid grid-4">
        {_kpi_tile("Active 30d", f"{_e(ex.get('active_30d'))} / {_e(ex.get('licensed_seats'))}")}
        {_kpi_tile("Logins", _e(ex.get('logins')))}
        {_kpi_tile("Hours", _e(ex.get('hours')))}
        {_kpi_tile("Abi queries", _e(ex.get('abi_queries')))}
      </div>"""

    # Slide 5 — Contract Summary
    modules_html = ", ".join(_e(m) for m in (cs.get("modules") or [])) or "—"
    s5_body = f"""
      <div class="grid grid-3">
        {_kpi_tile("Contract start", _e(cs.get('start')))}
        {_kpi_tile("Contract end / renewal", _e(cs.get('end')), f"{_e(cs.get('days_to_renewal'))} days" if cs.get('days_to_renewal') is not None else "")}
        {_kpi_tile("Term", _e(cs.get('term')))}
        {_kpi_tile("ACV", _money(cs.get('acv_usd')))}
        {_kpi_tile("Billing", _e(cs.get('billing')))}
        {_kpi_tile("Licensed seats", _e(cs.get('seats_current')), f"+{_e(cs.get('seats_proposed'))} in proposal" if cs.get('seats_proposed') else "")}
        {_kpi_tile("Geography", _e(cs.get('geography')))}
        {_kpi_tile("Sponsor", _e(cs.get('sponsor_name')), _e(cs.get('sponsor_role')))}
      </div>
      <h3 class="section">Modules in scope</h3>
      <p style="font-size:12px;color:#5a7896">{modules_html}</p>"""

    # Slide 6 — Accomplishments
    acc_cards = "".join(
        f'<div class="dollar-card"><div class="amt">{_money(a.get("amount_usd"))}</div>'
        f'<div class="label">{_e(a.get("title"))}</div>'
        f'<div class="desc">{_e(a.get("bu"))} · {_e(a.get("description"))}</div></div>'
        for a in acc
    ) or '<div class="dollar-card">—</div>'
    s6_body = acc_cards

    # Slide 7 — Upsell Pipeline
    up_rows = "".join(
        f'<tr><td>{_e(p.get("title"))}</td>'
        f'<td><span class="pill blue">{_e(p.get("status"))}</span></td>'
        f'<td>{_money(p.get("value_usd"))}</td>'
        f'<td>{_e(p.get("prob"))}%</td>'
        f'<td>{_e(p.get("when_text"))}</td></tr>'
        for p in up
    ) or '<tr><td colspan="5">—</td></tr>'
    s7_body = f"""
      <table>
        <thead><tr><th>Conversation</th><th>Status</th><th>Value</th><th>Prob</th><th>When</th></tr></thead>
        <tbody>{up_rows}</tbody>
      </table>"""

    # Slide 8 — Scores · 12
    band_pill = {"HEALTHY": "green", "AT RISK": "amber", "CRITICAL": "red"}

    def _band_cell(band: Any) -> str:
        if not band:
            return "—"
        cls = band_pill.get(band, "amber")
        return f'<span class="pill {cls}">{_e(band)}</span>'

    score_rows = "".join(
        f'<tr><td>{_e(r.get("name"))}</td>'
        f'<td style="font-weight:800;font-size:16px">{_e(r.get("score"))}</td>'
        f'<td>{_band_cell(r.get("band"))}</td></tr>'
        for r in sc
    )
    s8_body = f"""
      <table>
        <thead><tr><th>Dimension</th><th>Score</th><th>Band</th></tr></thead>
        <tbody>{score_rows}</tbody>
      </table>"""

    # Slide 9 — Subscribers & Engagement (12-month login chart)
    chart9_url = _png_data_url(charts.get("slide9", b""))
    chart9_html = (
        f'<img src="{chart9_url}" alt="12-month login trend" '
        f'style="width:100%;max-width:1000px;display:block;margin:20px auto 0">'
        if chart9_url else
        '<div style="padding:32px;text-align:center;color:#8496b0;font-size:11px;'
        'border:1px dashed #c5d0e0;border-radius:8px;margin-top:20px">'
        'Not enough monthly data yet — chart will appear once usage telemetry '
        'has ≥2 months of history.</div>'
    )
    s9_body = f"""
      <div class="grid grid-4">
        {_kpi_tile("LICENSED SEATS", _e(se.get('licensed_seats')), f"+{_e(se.get('seats_proposed'))} in proposal" if se.get('seats_proposed') else "")}
        {_kpi_tile("ACTIVE 30d", _e(se.get('active_30d')), f"{_e(se.get('activation_pct'))}% activation" if se.get('activation_pct') is not None else "")}
        {_kpi_tile("Logins (period)", _e(se.get('logins_total')))}
        {_kpi_tile("Hours (period)", _e(se.get('hours_total')))}
      </div>
      {chart9_html}"""

    # Slide 10 — Live.ai Category Watch
    _heat_pill = {"hot": "red", "warm": "amber", "whitespace": "blue", "cold": "green"}

    def _heat_cell(h: Any) -> str:
        cls = _heat_pill.get(str(h or "").lower(), "amber")
        return f'<span class="pill {cls}">{_e(h)}</span>'

    cat_rows = "".join(
        f'<tr><td>{_e(c.get("name"))}</td>'
        f'<td>{_e(c.get("visits"))}</td>'
        f'<td>{_heat_cell(c.get("heat"))}</td></tr>'
        for c in (la.get("top_cats") or [])
    ) or '<tr><td colspan="3">—</td></tr>'
    chart10_url = _png_data_url(charts.get("slide10", b""))
    chart10_html = (
        f'<img src="{chart10_url}" alt="Top categories by visits" '
        f'style="width:100%;max-width:1000px;display:block;margin:14px auto 0">'
        if chart10_url else ""
    )
    s10_body = f"""
      <div class="grid grid-4">
        {_kpi_tile("Subscribers", f"{_e(la.get('subscribers'))} of {_e(la.get('total_subs'))}")}
        {_kpi_tile("Categories unlocked", _e(la.get('categories_unlocked')),
                   f"{_e(la.get('ent_cats'))} Ent + {_e(la.get('non_ent_cats'))} Non-Ent" if la.get('ent_cats') is not None else "")}
        {_kpi_tile("Avg cat/user", _e(la.get('avg_per_user')), f"benchmark {_e(la.get('benchmark'))}")}
        {_kpi_tile("Top heat", _e((la.get('top_cats') or [{}])[0].get('name') if la.get('top_cats') else None))}
      </div>
      <h3 class="section">Top categories</h3>
      <table><thead><tr><th>Category</th><th>Visits</th><th>Heat</th></tr></thead>
        <tbody>{cat_rows}</tbody></table>
      {chart10_html}"""

    # Slide 11 — Inflation Watch GIT (monthly trend line)
    chart11_url = _png_data_url(charts.get("slide11", b""))
    chart11_html = (
        f'<img src="{chart11_url}" alt="Inflation Watch monthly trend" '
        f'style="width:100%;max-width:1000px;display:block;margin:20px auto 0">'
        if chart11_url else
        '<div style="padding:32px;text-align:center;color:#8496b0;font-size:11px;'
        'border:1px dashed #c5d0e0;border-radius:8px;margin-top:20px">'
        'No trend data yet — chart will appear once inflation telemetry has '
        '≥2 months of history.</div>'
    )
    s11_body = f"""
      <div class="grid grid-3">
        {_kpi_tile("Categories tracked", f"{_e(iw.get('categories_tracked'))} of {_e(iw.get('categories_in_scope'))} in scope")}
        {_kpi_tile("Views (period)", _e(iw.get('views_period')))}
        {_kpi_tile("Negotiation prep runs", _e(iw.get('neg_prep_runs')))}
      </div>
      {chart11_html}"""

    # Slide 12 — Closer
    s12 = f"""<section class="slide cover">
      <h1>Thank you, {_e(cover.get('account_name'))}</h1>
      <div class="meta">Next review scheduled per cadence</div>
      <div class="title">{_e(cover.get('title'))}</div>
      <div class="brand">{_e(cover.get('footer'))}</div>
    </section>"""

    body = "".join([
        s1,
        slide("CUSTOM CREDITS", "Custom Credits", s2_body),
        slide("RISKS · OPEN · ASKS", "Risks · Open · Asks", s3_body),
        slide("EXECUTIVE SNAPSHOT", "Executive Snapshot", s4_body),
        slide("CONTRACT SUMMARY", "Contract Summary", s5_body),
        slide("ACCOMPLISHMENTS & MILESTONES", "Accomplishments & Milestones", s6_body),
        slide("UPSELL & EXPANSION PIPELINE", "Upsell & Expansion Pipeline", s7_body),
        slide("ACCOUNT HEALTH · 12 SCORES", "Account Health · 12 scores", s8_body),
        slide("SUBSCRIBERS & ENGAGEMENT", "Subscribers & Engagement", s9_body),
        slide("LIVE.AI · CATEGORY WATCH", "Live.ai · Category Watch", s10_body),
        slide("INFLATION WATCH GIT", "Inflation Watch GIT", s11_body),
        s12,
    ])

    title_text = _html.escape(f"{cover.get('account_name', 'Business Review')} · {cover.get('period_label', '')}")
    return f"""<!doctype html>
<html><head><meta charset="utf-8"><title>{title_text}</title>
<style>{_CSS}</style></head>
<body>{body}</body></html>"""


# ============================================================
# PDF renderer — weasyprint HTML→PDF
# ============================================================


def render_pdf(html: str) -> bytes:
    """Render the HTML doc to PDF bytes.

    Returns empty bytes when:
      * weasyprint isn't installed
      * native libs (pango/cairo) are missing — common on a clean macOS dev
        box (fix: `brew install pango`); on Render Linux they ship via apt
    The route handler treats empty bytes as "PDF not available" (HTTP 409
    on /download?format=pdf), so a missing dep doesn't crash generation.
    """
    try:
        from weasyprint import HTML  # type: ignore
    except ImportError:
        logger.warning("weasyprint not installed; returning empty PDF")
        return b""
    except OSError as e:
        logger.warning("weasyprint native libs missing (%s); returning empty PDF", e)
        return b""
    try:
        buf = BytesIO()
        HTML(string=html).write_pdf(buf)
        return buf.getvalue()
    except Exception as e:  # noqa: BLE001
        logger.exception("weasyprint render failed: %s", e)
        return b""


# ============================================================
# PPTX renderer — python-pptx 12 slides
# ============================================================


def render_pptx(snapshot: dict[str, Any], *, charts: dict[str, bytes] | None = None) -> bytes:
    try:
        from pptx import Presentation  # type: ignore
        from pptx.dml.color import RGBColor  # type: ignore
        from pptx.util import Inches, Pt  # type: ignore
    except ImportError:
        logger.warning("python-pptx not installed; returning empty PPTX")
        return b""
    if charts is None:
        charts = _build_chart_pngs(snapshot)

    s = snapshot
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    NAVY = RGBColor(0x00, 0x11, 0x37)
    INDIGO = RGBColor(0x4A, 0x00, 0xF8)
    BUMBLEBEE = RGBColor(0xFF, 0xE6, 0x1E)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    MUTED = RGBColor(0x5A, 0x78, 0x96)
    DARK = RGBColor(0x0D, 0x1B, 0x2E)

    def _add_text(slide, x: float, y: float, w: float, h: float, text: str, *,
                  size: int = 14, bold: bool = False, color: RGBColor = DARK,
                  align: str = "left") -> None:
        tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        if align == "center":
            from pptx.enum.text import PP_ALIGN
            p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Manrope"

    def _add_cover(title_top: str, big_title: str, accent: str, footer: str) -> None:
        slide = prs.slides.add_slide(blank)
        bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = NAVY
        bg.line.fill.background()
        _add_text(slide, 1, 2.2, 11.3, 0.6, title_top, size=16, color=RGBColor(0xA3, 0xB6, 0xE0), align="center")
        _add_text(slide, 1, 2.9, 11.3, 1.2, big_title, size=42, bold=True, color=WHITE, align="center")
        _add_text(slide, 1, 4.4, 11.3, 0.6, accent, size=20, bold=True, color=BUMBLEBEE, align="center")
        _add_text(slide, 1, 6.6, 11.3, 0.4, footer, size=12, color=RGBColor(0xA3, 0xB6, 0xE0), align="center")

    def _add_content(eyebrow: str, title: str) -> Any:
        slide = prs.slides.add_slide(blank)
        _add_text(slide, 0.5, 0.4, 12.3, 0.3, "·  " + eyebrow.upper(), size=10, bold=True, color=MUTED)
        _add_text(slide, 0.5, 0.75, 12.3, 0.6, title, size=24, bold=True, color=NAVY)
        # underline
        ln = slide.shapes.add_shape(1, Inches(0.5), Inches(1.4), Inches(2), Inches(0.03))
        ln.fill.solid()
        ln.fill.fore_color.rgb = INDIGO
        ln.line.fill.background()
        return slide

    def _add_kpi(slide, x: float, y: float, w: float, label: str, value: str, note: str = "") -> None:
        tile = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(1.2))
        tile.fill.solid()
        tile.fill.fore_color.rgb = RGBColor(0xF7, 0xF9, 0xFD)
        tile.line.color.rgb = RGBColor(0xE4, 0xEA, 0xF6)
        _add_text(slide, x + 0.15, y + 0.1, w - 0.3, 0.25, label.upper(), size=8, bold=True, color=MUTED)
        _add_text(slide, x + 0.15, y + 0.4, w - 0.3, 0.5, value, size=20, bold=True, color=NAVY)
        if note:
            _add_text(slide, x + 0.15, y + 0.95, w - 0.3, 0.25, note, size=9, color=MUTED)

    def _fmt(v: Any) -> str:
        if v is None or v == "":
            return "—"
        return str(v)

    cover = s["cover"]
    cc = s["custom_credits"]
    ro = s["risks_open_asks"]
    ex = s["executive_snapshot"]
    cs = s["contract_summary"]
    acc = s.get("accomplishments") or []
    up = s.get("upsell_pipeline") or []
    sc = s.get("scores_12") or []
    se = s["subscribers_engagement"]
    la = s["live_ai"]
    iw = s["inflation_watch"]

    # 1 Cover
    _add_cover(
        cover.get("period_label", ""),
        cover.get("account_name", ""),
        cover.get("title", ""),
        cover.get("footer", "Beroe"),
    )

    # 2 Custom Credits
    sl = _add_content("Custom Credits", "Custom Credits")
    tiles = [
        ("FTE allocation", _fmt(cc.get("fte"))),
        ("Hours purchased", _fmt(cc.get("hours_purchased"))),
        ("Hours consumed YTD", _fmt(cc.get("hours_consumed"))),
        ("Hours remaining", _fmt(cc.get("hours_remaining"))),
        ("Fixed-fee instances", f"{_fmt(cc.get('fixed_fee_done'))}/{_fmt(cc.get('fixed_fee_total'))}"),
        ("Instances in flight", _fmt(cc.get("instances_in_flight"))),
        ("Infinity Slots active", _fmt(cc.get("slots_active"))),
        ("Slot refresh cadence", _fmt(cc.get("slot_cadence"))),
    ]
    for i, (lab, val) in enumerate(tiles):
        x = 0.5 + (i % 4) * 3.1
        y = 1.7 + (i // 4) * 1.4
        _add_kpi(sl, x, y, 3.0, lab, val)

    # 3 Risks · Open · Asks
    sl = _add_content("Risks · Open · Asks", "Risks · Open · Asks")
    for col, (label, items) in enumerate([
        ("Open items (Beroe-side)", ro.get("open_items") or ["—"]),
        ("Risks", ro.get("risks") or ["—"]),
        ("Asks from client", ro.get("asks") or ["—"]),
    ]):
        x = 0.5 + col * 4.3
        _add_text(sl, x, 1.7, 4.0, 0.35, label, size=12, bold=True, color=NAVY)
        for i, txt in enumerate(items[:8]):
            _add_text(sl, x, 2.1 + i * 0.45, 4.0, 0.45, "• " + str(txt), size=11, color=DARK)

    # 4 Executive Snapshot
    sl = _add_content("Executive Snapshot", "Executive Snapshot")
    _add_kpi(sl, 0.5, 1.7, 3.0, "Account Health", f"{_fmt(ex.get('health'))} / 100",
             _fmt(ex.get("health_trend") or ex.get("health_band") or ""))
    _add_kpi(sl, 3.7, 1.7, 3.0, "ARR in scope", _money(ex.get("arr_usd")),
             _fmt(ex.get("expansion_note") or ""))
    _add_kpi(sl, 6.9, 1.7, 3.0, "Cost avoidance", _money(ex.get("cost_avoidance_usd")),
             f"target {_money(ex.get('cost_target_usd'))}" if ex.get("cost_target_usd") else "")
    _add_kpi(sl, 10.1, 1.7, 2.7, "Renewal", _fmt(ex.get("renewal_date")),
             f"{ex.get('days_to_renewal')} days" if ex.get("days_to_renewal") is not None else "")
    _add_text(sl, 0.5, 3.3, 5, 0.3, "Engagement", size=12, bold=True, color=NAVY)
    _add_kpi(sl, 0.5, 3.7, 3.0, "Active 30d",
             f"{_fmt(ex.get('active_30d'))} / {_fmt(ex.get('licensed_seats'))}")
    _add_kpi(sl, 3.7, 3.7, 3.0, "Logins", _fmt(ex.get("logins")))
    _add_kpi(sl, 6.9, 3.7, 3.0, "Hours", _fmt(ex.get("hours")))
    _add_kpi(sl, 10.1, 3.7, 2.7, "Abi queries", _fmt(ex.get("abi_queries")))

    # 5 Contract Summary
    sl = _add_content("Contract Summary", "Contract Summary")
    items = [
        ("Contract start", _fmt(cs.get("start"))),
        ("Contract end / renewal", _fmt(cs.get("end"))),
        ("Term", _fmt(cs.get("term"))),
        ("ACV", _money(cs.get("acv_usd"))),
        ("Billing", _fmt(cs.get("billing"))),
        ("Licensed seats", _fmt(cs.get("seats_current"))),
        ("Geography", _fmt(cs.get("geography"))),
        ("Sponsor", f"{_fmt(cs.get('sponsor_name'))} · {_fmt(cs.get('sponsor_role'))}"),
    ]
    for i, (lab, val) in enumerate(items):
        x = 0.5 + (i % 4) * 3.1
        y = 1.7 + (i // 4) * 1.4
        _add_kpi(sl, x, y, 3.0, lab, val)
    mods = ", ".join(cs.get("modules") or []) or "—"
    _add_text(sl, 0.5, 4.7, 12.3, 0.3, "Modules in scope", size=12, bold=True, color=NAVY)
    _add_text(sl, 0.5, 5.0, 12.3, 1.5, mods, size=11, color=DARK)

    # 6 Accomplishments
    sl = _add_content("Accomplishments & Milestones", "Accomplishments & Milestones")
    for i, a in enumerate((acc or [{}])[:6]):
        x = 0.5 + (i % 3) * 4.2
        y = 1.7 + (i // 3) * 2.5
        _add_text(sl, x, y, 4.0, 0.5, _money(a.get("amount_usd")), size=22, bold=True, color=INDIGO)
        _add_text(sl, x, y + 0.6, 4.0, 0.35, _fmt(a.get("title")), size=11, bold=True, color=NAVY)
        _add_text(sl, x, y + 1.0, 4.0, 0.3, _fmt(a.get("bu")), size=9, color=MUTED)
        _add_text(sl, x, y + 1.3, 4.0, 1.0, _fmt(a.get("description")), size=10, color=DARK)

    # 7 Upsell pipeline
    sl = _add_content("Upsell & Expansion Pipeline", "Upsell & Expansion Pipeline")
    _add_text(sl, 0.5, 1.7, 5, 0.3, "Conversation", size=10, bold=True, color=MUTED)
    _add_text(sl, 5.0, 1.7, 2, 0.3, "Status", size=10, bold=True, color=MUTED)
    _add_text(sl, 7.0, 1.7, 2, 0.3, "Value", size=10, bold=True, color=MUTED)
    _add_text(sl, 9.0, 1.7, 1.5, 0.3, "Prob", size=10, bold=True, color=MUTED)
    _add_text(sl, 10.5, 1.7, 2.5, 0.3, "When", size=10, bold=True, color=MUTED)
    for i, p in enumerate(up[:8]):
        y = 2.1 + i * 0.5
        _add_text(sl, 0.5, y, 4.5, 0.4, _fmt(p.get("title")), size=11, color=DARK)
        _add_text(sl, 5.0, y, 2, 0.4, _fmt(p.get("status")), size=10, color=INDIGO, bold=True)
        _add_text(sl, 7.0, y, 2, 0.4, _money(p.get("value_usd")), size=11, color=DARK, bold=True)
        _add_text(sl, 9.0, y, 1.5, 0.4, f"{p.get('prob', 0)}%", size=11, color=DARK)
        _add_text(sl, 10.5, y, 2.5, 0.4, _fmt(p.get("when_text")), size=10, color=MUTED)

    # 8 Scores · 12
    sl = _add_content("Account Health · 12 scores", "Account Health · 12 scores")
    for i, r in enumerate(sc):
        x = 0.5 + (i % 4) * 3.1
        y = 1.7 + (i // 4) * 1.5
        _add_kpi(sl, x, y, 3.0, r.get("name", ""), str(_fmt(r.get("score"))), _fmt(r.get("band")))

    # 9 Subscribers & Engagement
    sl = _add_content("Subscribers & Engagement", "Subscribers & Engagement")
    _add_kpi(sl, 0.5, 1.7, 3.0, "Licensed seats", _fmt(se.get("licensed_seats")),
             f"+{se.get('seats_proposed')} in proposal" if se.get("seats_proposed") else "")
    _add_kpi(sl, 3.7, 1.7, 3.0, "Active 30d", _fmt(se.get("active_30d")),
             f"{se.get('activation_pct')}% activation" if se.get("activation_pct") is not None else "")
    _add_kpi(sl, 6.9, 1.7, 3.0, "Logins (period)", _fmt(se.get("logins_total")))
    _add_kpi(sl, 10.1, 1.7, 2.7, "Hours (period)", _fmt(se.get("hours_total")))
    chart9_bytes = charts.get("slide9", b"")
    if chart9_bytes:
        sl.shapes.add_picture(BytesIO(chart9_bytes), Inches(0.5), Inches(3.2),
                              width=Inches(12.3), height=Inches(3.8))
    else:
        _add_text(sl, 0.5, 3.5, 12.3, 0.4,
                  "Trend chart unavailable — need ≥2 months of usage data.",
                  size=11, color=MUTED, align="center")

    # 10 Live.ai
    sl = _add_content("Live.ai · Category Watch", "Live.ai · Category Watch")
    _add_kpi(sl, 0.5, 1.7, 3.0, "Subscribers",
             f"{_fmt(la.get('subscribers'))} of {_fmt(la.get('total_subs'))}")
    _add_kpi(sl, 3.7, 1.7, 3.0, "Categories unlocked", _fmt(la.get("categories_unlocked")),
             f"{la.get('ent_cats') or '—'} Ent + {la.get('non_ent_cats') or '—'} Non-Ent")
    _add_kpi(sl, 6.9, 1.7, 3.0, "Avg cat/user", _fmt(la.get("avg_per_user")),
             f"benchmark {la.get('benchmark')}" if la.get("benchmark") else "")
    chart10_bytes = charts.get("slide10", b"")
    if chart10_bytes:
        sl.shapes.add_picture(BytesIO(chart10_bytes), Inches(0.5), Inches(3.2),
                              width=Inches(12.3), height=Inches(3.8))
    else:
        _add_text(sl, 0.5, 3.5, 12.3, 0.4,
                  "No top-category data — populate platform_intel.cat_intel.top_cats.",
                  size=11, color=MUTED, align="center")

    # 11 Inflation Watch GIT
    sl = _add_content("Inflation Watch GIT", "Inflation Watch GIT")
    _add_kpi(sl, 0.5, 1.7, 4.0, "Categories tracked",
             f"{_fmt(iw.get('categories_tracked'))} of {_fmt(iw.get('categories_in_scope'))}")
    _add_kpi(sl, 4.7, 1.7, 4.0, "Views (period)", _fmt(iw.get("views_period")))
    _add_kpi(sl, 8.9, 1.7, 3.9, "Negotiation prep runs", _fmt(iw.get("neg_prep_runs")))
    chart11_bytes = charts.get("slide11", b"")
    if chart11_bytes:
        sl.shapes.add_picture(BytesIO(chart11_bytes), Inches(0.5), Inches(3.2),
                              width=Inches(12.3), height=Inches(3.8))
    else:
        _add_text(sl, 0.5, 3.5, 12.3, 0.4,
                  "No inflation trend data — populate platform_intel.inflation_watch.trend_monthly.",
                  size=11, color=MUTED, align="center")

    # 12 Closer
    _add_cover(
        "Next review scheduled per cadence",
        f"Thank you, {cover.get('account_name', '')}",
        cover.get("title", ""),
        cover.get("footer", "Beroe"),
    )

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
